import urllib.request
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_khushbu_cyber_tools_perfected():
    prs = Presentation()
    
    # Theme: Deep Forest & Cyber Lime
    DARK_BG = RGBColor(10, 20, 15)      
    ACCENT_LIME = RGBColor(163, 230, 53) 
    TEXT_WHITE = RGBColor(243, 244, 246)

    # 3D SECURITY ICON SET
    icons = {
        "title": "https://cdn-icons-png.flaticon.com/512/2092/2092663.png",
        "intro": "https://cdn-icons-png.flaticon.com/512/1021/1021656.png",
        "importance": "https://cdn-icons-png.flaticon.com/512/1162/1162939.png",
        "types": "https://cdn-icons-png.flaticon.com/512/2438/2438078.png",
        "network": "https://cdn-icons-png.flaticon.com/512/2885/2885417.png",
        "endpoint": "https://cdn-icons-png.flaticon.com/512/3067/3067451.png",
        "app": "https://cdn-icons-png.flaticon.com/512/1055/1055687.png",
        "popular": "https://cdn-icons-png.flaticon.com/512/1149/1149168.png",
        "emerging": "https://cdn-icons-png.flaticon.com/512/2103/2103633.png",
        "thanks": "https://cdn-icons-png.flaticon.com/512/14441/14441394.png"
    }

    def download_image(url):
        try:
            req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req, timeout=10) as response:
                return io.BytesIO(response.read())
        except: return None

    def apply_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    # 1. TITLE SLIDE - FIXED LAYOUT
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    apply_bg(slide)
    
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5.5), Inches(3))
    p1 = t_box.text_frame.add_paragraph()
    p1.text = "CYBER SECURITY\nTOOLS & TECH"
    p1.font.size, p1.font.bold, p1.font.color.rgb = Pt(48), True, ACCENT_LIME
    
    p2 = t_box.text_frame.add_paragraph()
    p2.text = "Protecting Digital Systems and Data"
    p2.font.size, p2.font.color.rgb = Pt(22), TEXT_WHITE

    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(5), Inches(1))
    ip = info_box.text_frame.add_paragraph()
    ip.text = "Name: Khushbu Patel\nEnroll No: 23ss02it137"
    ip.font.size, ip.font.color.rgb = Pt(18), TEXT_WHITE
    
    img = download_image(icons["title"])
    if img: slide.shapes.add_picture(img, Inches(6.2), Inches(1.8), height=Inches(3.5))

    # --- CONTENT SLIDES - BALANCED SPLIT LAYOUT ---
    content_map = [
        ("Introduction", ["Protects systems, networks, and data.", "Helps detect, prevent, and respond.", "Rising importance due to cybercrime."], "intro"),
        ("Why Tools are Important", ["Protect sensitive data & privacy.", "Prevent hacking and malware.", "Ensure system reliability and trust."], "importance"),
        ("Types of Security Tools", ["Network & Endpoint Security.", "Application & Cloud Security.", "Identity & Access Management."], "types"),
        ("Network Security Tools", ["Firewalls & VPNs.", "Intrusion Detection (IDS).", "Intrusion Prevention (IPS)."], "network"),
        ("Endpoint Security Tools", ["Antivirus & Anti-malware.", "Device encryption methods.", "Mobile security applications."], "endpoint"),
        ("Application Security", ["Web App Firewalls (WAF).", "Code analysis & Bug tracking.", "Penetration testing tools."], "app"),
        ("Popular Security Tools", ["Wireshark & Nmap.", "Metasploit & Kali Linux.", "Norton & McAfee antivirus."], "popular"),
        ("Emerging Technologies", ["AI and Machine Learning.", "Blockchain & Zero Trust.", "Cloud security tech."], "emerging"),
        ("Conclusion", ["Essential for modern safety.", "Requires continuous updates.", "Awareness is key to protection."], "emerging")
    ]

    for title_text, points, icon_key in content_map:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(slide)
        
        # Title - Top Left
        t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tp = t_box.text_frame.add_paragraph()
        tp.text = title_text
        tp.font.size, tp.font.bold, tp.font.color.rgb = Pt(36), True, ACCENT_LIME
        
        # Content - Left Split (Reduced width to 6.2 inches)
        c_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6.2), Inches(4.5))
        ctf = c_box.text_frame
        ctf.word_wrap = True
        for point in points:
            p = ctf.add_paragraph()
            p.text = f"• {point}"
            p.font.color.rgb, p.font.size, p.space_after = TEXT_WHITE, Pt(20), Pt(18)
        
        # Image - Right Split (Positioned at 6.8 inches)
        img = download_image(icons[icon_key])
        if img: slide.shapes.add_picture(img, Inches(6.8), Inches(2.2), width=Inches(2.8))

    # FINAL THANK YOU SLIDE
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide)
    tx_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(10), Inches(2))
    p_main = tx_box.text_frame.add_paragraph()
    p_main.text = "Thank You"
    p_main.font.size, p_main.font.bold, p_main.font.color.rgb = Pt(64), True, ACCENT_LIME
    p_main.alignment = PP_ALIGN.CENTER

    p_info = tx_box.text_frame.add_paragraph()
    p_info.text = "Any Questions?\nKhushbu Patel | 23ss02it137"
    p_info.font.size, p_info.font.color.rgb = Pt(22), TEXT_WHITE
    p_info.alignment = PP_ALIGN.CENTER

    prs.save("N:\\projects\\project\\sem6\\pawX\\cyber_tools_khushbu_v2.pptx")
    print("Success! Balanced layout created for Khushbu.")

if __name__ == "__main__":
    create_khushbu_cyber_tools_perfected()